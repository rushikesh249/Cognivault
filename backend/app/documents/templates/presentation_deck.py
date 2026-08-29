"""PPTX Presentation Deck Generator (TRD Section 22, Component #17)."""

import datetime
from pathlib import Path
from typing import Any, Dict, List, Tuple
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Palette
COLOR_NAVY = RGBColor(16, 44, 87)       # Primary Header #102C57
COLOR_ICE = RGBColor(240, 244, 248)     # Background Box #F0F4F8
COLOR_CARD_BG = RGBColor(249, 251, 252) # Card Fill #F9FBFC
COLOR_DARK_TEXT = RGBColor(33, 37, 41)  # Charcoal Text
COLOR_MUTED = RGBColor(108, 117, 125)   # Muted Gray
COLOR_ALERT = RGBColor(180, 0, 0)       # Alert Red #B40000
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(208, 215, 222)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"
DISCLAIMER_TEXT = "AI-Generated Draft — For Internal Engineering Review Only (Non-Certified Verdict)"


def add_slide_header_and_footer(slide, title_text: str, category: str = "EXECUTIVE ENGINEERING BRIEFING"):
    """Render consistent 16:9 sovereign banner, title, and bottom disclaimer."""
    # Top Category Tracker
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = FONT_HEADING
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_MUTED

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_NAVY

    # Decorative Line Under Title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_NAVY
    line.line.color.rgb = COLOR_NAVY

    # Footer Disclaimer Banner
    foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    tf_foot = foot_box.text_frame
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = DISCLAIMER_TEXT
    p_foot.font.name = FONT_BODY
    p_foot.font.size = Pt(8.5)
    p_foot.font.italic = True
    p_foot.font.color.rgb = COLOR_ALERT
    p_foot.alignment = PP_ALIGN.CENTER


def render_presentation_deck(data: Dict[str, Any], output_path: Path) -> Path:
    """
    Render 5-slide 16:9 executive presentation deck.

    Slides:
    1. Title & Executive Metadata
    2. Inspection Overview & Scope
    3. Critical Findings & Observations
    4. Compliance Gaps & SOP Citations
    5. Actionable Recommendations & Sign-off
    """
    prs = pptx.Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Completely blank slide layout

    title_main = str(data.get("title", "EQUIPMENT INSPECTION & COMPLIANCE EVALUATION"))
    task_id = str(data.get("task_id", "TASK-AUTONOMOUS-01"))
    facility = str(data.get("facility", "Primary Refining Unit & Flare Header"))
    timestamp = str(data.get("timestamp", datetime.datetime.now(datetime.timezone.utc).strftime("%Y-%m-%d %H:%M UTC")))
    status = str(data.get("status", "ACTION REQUIRED — NON-COMPLIANCE DETECTED"))

    # =========================================================================
    # Slide 1: Title & Executive Metadata
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)

    # Top Brand Bar
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_NAVY
    top_bar.line.fill.background()

    # Title Box
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.8))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "SOVEREIGN AGENTIC AI WORKBENCH"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_MUTED

    p2 = tf1.add_paragraph()
    p2.text = title_main.upper()
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_NAVY
    p2.space_before = Pt(8)

    # Metadata Grid Card
    meta_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.2), Inches(11.333), Inches(2.6))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = COLOR_ICE
    meta_card.line.color.rgb = COLOR_BORDER

    tf_m = meta_card.text_frame
    tf_m.word_wrap = True
    tf_m.vertical_anchor = MSO_ANCHOR.MIDDLE

    meta_items = [
        ("Task Reference ID:", task_id),
        ("Facility / Asset Unit:", facility),
        ("Evaluation Timestamp:", timestamp),
        ("Compliance Status:", status),
    ]

    for idx, (label, val) in enumerate(meta_items):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        run_l = p.add_run()
        run_l.text = f"{label:<26} "
        run_l.font.name = FONT_BODY
        run_l.font.size = Pt(11.5)
        run_l.font.bold = True
        run_l.font.color.rgb = COLOR_NAVY

        run_v = p.add_run()
        run_v.text = f"{val}\n"
        run_v.font.name = FONT_BODY
        run_v.font.size = Pt(11.5)
        if "ACTION REQUIRED" in val:
            run_v.font.bold = True
            run_v.font.color.rgb = COLOR_ALERT
        else:
            run_v.font.color.rgb = COLOR_DARK_TEXT
        p.space_after = Pt(6)

    # Footer
    add_slide_header_and_footer(s1, "", category="SOVEREIGN ON-PREMISE AI")

    # =========================================================================
    # Slide 2: Inspection Overview & Scope
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(s2, "1. Inspection Overview & Analysis Scope")

    summary_text = str(data.get(
        "summary",
        "Autonomous document intelligence analysis was executed on the submitted technical inspection report. "
        "The document was analyzed locally with zero external network egress, extracting structural defect measurements "
        "and evaluating compliance against indexed safety standards and operating procedures."
    ))

    # Executive Overview Box
    card_ov = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2))
    card_ov.fill.solid()
    card_ov.fill.fore_color.rgb = COLOR_CARD_BG
    card_ov.line.color.rgb = COLOR_BORDER

    tf_ov = card_ov.text_frame
    tf_ov.word_wrap = True
    p_ov_h = tf_ov.paragraphs[0]
    p_ov_h.text = "Executive Analysis Summary"
    p_ov_h.font.name = FONT_HEADING
    p_ov_h.font.size = Pt(13)
    p_ov_h.font.bold = True
    p_ov_h.font.color.rgb = COLOR_NAVY

    p_ov_b = tf_ov.add_paragraph()
    p_ov_b.text = summary_text
    p_ov_b.font.name = FONT_BODY
    p_ov_b.font.size = Pt(11)
    p_ov_b.font.color.rgb = COLOR_DARK_TEXT
    p_ov_b.space_before = Pt(8)

    # Scope & Methodology 3 Pillars
    pillars = [
        ("OCR & Document Ingestion", "Local OCR engine processed raw technical scans, extracting tables, text blocks, and defect notes."),
        ("RAG Semantic Indexing", "Discovered defects matched against local SOP embeddings with 0.55 similarity threshold."),
        ("Compliance Rule Engine", "Defect parameters evaluated against statutory safety and maintenance tolerances."),
    ]
    for idx, (p_title, p_desc) in enumerate(pillars):
        x = Inches(0.8 + idx * 4.0)
        p_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.1), Inches(3.7), Inches(2.4))
        p_card.fill.solid()
        p_card.fill.fore_color.rgb = COLOR_ICE
        p_card.line.color.rgb = COLOR_BORDER
        tf_p = p_card.text_frame
        tf_p.word_wrap = True

        p_h = tf_p.paragraphs[0]
        p_h.text = f"{idx+1}. {p_title}"
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_NAVY

        p_b = tf_p.add_paragraph()
        p_b.text = p_desc
        p_b.font.name = FONT_BODY
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = COLOR_DARK_TEXT
        p_b.space_before = Pt(6)

    # =========================================================================
    # Slide 3: Critical Findings & Observations
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(s3, "2. Critical Inspection Findings & Anomalies")

    raw_findings = data.get("critical_findings", [
        "Corrosion fatigue detected on primary discharge flange bolts exceeding 1.5mm wall thinning threshold.",
        "Pressure relief valve PRV-204 calibration interval exceeded maximum allowable 12-month limit.",
        "Emergency shutdown bypass valve seal integrity compromised with visible weeping of seal fluid."
    ])

    for idx, finding in enumerate(raw_findings[:3]):
        y = Inches(1.6 + idx * 1.65)
        f_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.4))
        f_card.fill.solid()
        f_card.fill.fore_color.rgb = COLOR_CARD_BG
        f_card.line.color.rgb = COLOR_BORDER

        tf_f = f_card.text_frame
        tf_f.word_wrap = True

        p_fh = tf_f.paragraphs[0]
        p_fh.text = f"FINDING #{idx+1:02d} — DEFECT ANOMALY DETECTED"
        p_fh.font.name = FONT_HEADING
        p_fh.font.size = Pt(11)
        p_fh.font.bold = True
        p_fh.font.color.rgb = COLOR_ALERT if idx == 0 else COLOR_NAVY

        p_fb = tf_f.add_paragraph()
        p_fb.text = str(finding)
        p_fb.font.name = FONT_BODY
        p_fb.font.size = Pt(10.5)
        p_fb.font.color.rgb = COLOR_DARK_TEXT
        p_fb.space_before = Pt(4)

    # =========================================================================
    # Slide 4: Compliance Gaps & SOP Citations
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(s4, "3. Compliance Gaps & Authoritative Standard Citations")

    raw_gaps = data.get("compliance_gaps", [
        ("Discharge Flange Wall Thinning", "Safety SOP - Section 4.2 Emergency Shutdown Systems (p.12)", "CRITICAL NON-COMPLIANCE"),
        ("PRV-204 Recertification Overdue", "Equipment Standards - Section 11.4 Relief Valve Recertification (p.56)", "MAJOR GAP"),
        ("Seal Integrity Weeping", "Maintenance Manual - Section 8.1 Flange Integrity & Bolt Torquing (p.34)", "MODERATE GAP"),
    ])

    rows = len(raw_gaps) + 1
    table_shape = s4.shapes.add_table(rows, 3, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(3.8)
    table.columns[1].width = Inches(5.4)
    table.columns[2].width = Inches(2.5)

    headers = ["Observed Finding / Defect", "Authoritative SOP Citation", "Compliance Rating"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    for r_idx, item in enumerate(raw_gaps, start=1):
        defect, cit, rating = item if isinstance(item, (list, tuple)) and len(item) >= 3 else (str(item), "Standard SOP", "NON-COMPLIANCE")
        vals = [defect, cit, rating]
        for c_idx, val in enumerate(vals):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG if r_idx % 2 == 0 else COLOR_ICE
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_ALERT if ("CRITICAL" in str(val) or "NON-COMPLIANCE" in str(val)) else COLOR_DARK_TEXT
            if c_idx == 2:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

    # =========================================================================
    # Slide 5: Actionable Recommendations & Sign-off
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(s5, "4. Actionable Engineering Recommendations & Sign-Off")

    raw_recs = data.get("recommendations", [
        "Initiate immediate scheduled depressurization and replacement of flange bolting assembly.",
        "Perform off-line hydrostatic bench testing and recertification for PRV-204 within 48 hours.",
        "Replace primary mechanical seal pack on pump P-102A prior to resuming continuous service."
    ])

    # Left Column: Recommendations (Width 7.2 in)
    rec_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.0))
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = COLOR_CARD_BG
    rec_box.line.color.rgb = COLOR_BORDER
    tf_r = rec_box.text_frame
    tf_r.word_wrap = True

    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Remediation Action Plan"
    p_rh.font.name = FONT_HEADING
    p_rh.font.size = Pt(13)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_NAVY

    for idx, rec in enumerate(raw_recs, start=1):
        p_item = tf_r.add_paragraph()
        run_n = p_item.add_run()
        run_n.text = f"Action 4.{idx}: "
        run_n.font.name = FONT_BODY
        run_n.font.size = Pt(10.5)
        run_n.font.bold = True
        run_n.font.color.rgb = COLOR_ALERT if idx == 1 else COLOR_NAVY

        run_t = p_item.add_run()
        run_t.text = str(rec)
        run_t.font.name = FONT_BODY
        run_t.font.size = Pt(10)
        run_t.font.color.rgb = COLOR_DARK_TEXT
        p_item.space_before = Pt(8)

    # Right Column: Sign-off card (Width 4.2 in)
    sign_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.6), Inches(4.2), Inches(5.0))
    sign_box.fill.solid()
    sign_box.fill.fore_color.rgb = COLOR_ICE
    sign_box.line.color.rgb = COLOR_BORDER
    tf_s = sign_box.text_frame
    tf_s.word_wrap = True

    p_sh = tf_s.paragraphs[0]
    p_sh.text = "Engineering Endorsement"
    p_sh.font.name = FONT_HEADING
    p_sh.font.size = Pt(13)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_NAVY

    p_prep = tf_s.add_paragraph()
    prep_status = data.get("status") or "Analyzed on-premise"
    p_prep.text = f"Prepared By:\nSovereign AI Workbench\n({prep_status})"
    p_prep.font.name = FONT_BODY
    p_prep.font.size = Pt(10)
    p_prep.font.color.rgb = COLOR_DARK_TEXT
    p_prep.space_before = Pt(14)

    p_rev = tf_s.add_paragraph()
    p_rev.text = "Reviewed & Endorsed By:\n\n____________________________\nLead Mechanical Maintenance Engineer\nLicense / Stamp: ______________"
    p_rev.font.name = FONT_BODY
    p_rev.font.size = Pt(10)
    p_rev.font.color.rgb = COLOR_DARK_TEXT
    p_rev.space_before = Pt(20)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
