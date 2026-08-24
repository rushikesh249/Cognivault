"""Unit tests for PPTX Presentation Deck Generator (TRD Section 22, Component #17)."""

from pathlib import Path
import pptx
import pytest

from backend.app.documents.templates.presentation_deck import render_presentation_deck


@pytest.fixture
def sample_presentation_data():
    return {
        "title": "EXECUTIVE MANAGEMENT COMPLIANCE REVIEW",
        "task_id": "TASK-PPTX-TEST-01",
        "facility": "Primary Refining Unit - Distillation Column",
        "timestamp": "2026-08-24 18:00 UTC",
        "status": "ACTION REQUIRED — NON-COMPLIANCE DETECTED",
        "summary": "Autonomous management briefing summarizing equipment inspection results, compliance gaps, and remediation actions.",
        "critical_findings": [
            "Refining column reflux line wall thinning below retirement thickness.",
            "Overhead condenser relief valve calibration overdue.",
            "High-pressure seal assembly weeping mineral seal oil.",
        ],
        "compliance_gaps": [
            ("Reflux Line Thinning", "Piping Inspection Code API 570 (p.24)", "CRITICAL NON-COMPLIANCE"),
            ("Relief Valve Overdue", "Safety SOP Section 4.2 (p.12)", "MAJOR GAP"),
            ("Seal Oil Weeping", "Pump Maintenance Standard Section 8.1 (p.45)", "MODERATE GAP"),
        ],
        "recommendations": [
            "Execute emergency clamp installation and schedule line replacement.",
            "Complete PRV bench test and certification within 48 hours.",
            "Replace seal cartridge during upcoming maintenance shift.",
        ],
    }


def test_pptx_generation_success(sample_presentation_data, tmp_path):
    """Verify 5-slide PPTX deck renders successfully with structured content."""
    output_path = tmp_path / "deck.pptx"
    res = render_presentation_deck(sample_presentation_data, output_path)

    assert res.exists()
    assert res.stat().st_size > 5000

    prs = pptx.Presentation(str(output_path))
    # Assert 5 slides generated
    assert len(prs.slides) == 5

    # Check Slide 1: Title & Metadata
    s1_text = "\n".join(shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame)
    assert "EXECUTIVE MANAGEMENT COMPLIANCE REVIEW" in s1_text
    assert "TASK-PPTX-TEST-01" in s1_text

    # Check Slide 2: Scope & Methodology
    s2_text = "\n".join(shape.text_frame.text for shape in prs.slides[1].shapes if shape.has_text_frame)
    assert "Inspection Overview" in s2_text

    # Check Slide 3: Findings
    s3_text = "\n".join(shape.text_frame.text for shape in prs.slides[2].shapes if shape.has_text_frame)
    assert "Critical Inspection Findings" in s3_text
    assert "wall thinning" in s3_text.lower()

    # Check Slide 4: Table with compliance gaps
    s4 = prs.slides[3]
    has_table = any(shape.has_table for shape in s4.shapes)
    assert has_table, "Slide 4 must contain a compliance gaps table"

    # Check Slide 5: Recommendations & Sign-off
    s5_text = "\n".join(shape.text_frame.text for shape in prs.slides[4].shapes if shape.has_text_frame)
    assert "Remediation Action Plan" in s5_text
    assert "Engineering Endorsement" in s5_text


def test_pptx_disclaimer_on_all_slides(sample_presentation_data, tmp_path):
    """Verify mandatory non-certified draft disclaimer is present across slides."""
    output_path = tmp_path / "deck_disclaimer.pptx"
    render_presentation_deck(sample_presentation_data, output_path)

    prs = pptx.Presentation(str(output_path))
    for idx, slide in enumerate(prs.slides):
        slide_text = "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        assert "AI-Generated Draft" in slide_text, f"Missing AI draft disclaimer on slide {idx+1}"
        assert "Non-Certified Verdict" in slide_text, f"Missing non-certified verdict notice on slide {idx+1}"
